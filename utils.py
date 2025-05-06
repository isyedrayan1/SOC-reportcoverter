from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd
from datetime import datetime

# Extract month and year from Date/Time column
def extract_month_year(df):
    try:
        # Parse the first Date/Time value
        date_str = df["Date/Time"].iloc[0]
        date_obj = datetime.strptime(date_str, "%m/%d/%Y %I:%M %p")
        month = date_obj.strftime("%B")
        year = date_obj.year
        return f"{month} 1st - {month} {date_obj.day}th, {year}"
    except Exception:
        return "Month X 2025"

# Process Excel data
def process_excel_data(file):
    # Read Excel file
    df = pd.read_excel(file)

    # Extract month/year
    month_year = extract_month_year(df)
    month = month_year.split(" ")[0]

    # Total alerts
    total_alerts = len(df)

    # Escalated alerts (Report Status = "Reported")
    escalated_alerts = len(df[df["Report Status"] == "Reported"])

    # Disposition breakdown
    disposition_counts = df["Disposition"].value_counts().to_dict()
    not_reported = disposition_counts.get("Not Reported", 0)
    reported = disposition_counts.get("Reported", 0)

    # Priority breakdown
    priority_counts = df["Priority"].value_counts().to_dict()
    priority_breakdown = (
        priority_counts.get("Low", 0),
        priority_counts.get("Medium", 0),
        priority_counts.get("High", 0)
    )

    # Previous month comparison (assuming placeholder data if not in Excel)
    previous_alerts = 383  # Placeholder (from March PPTX analysis)
    previous_escalated = 71  # Placeholder
    alerts_comparison = (previous_alerts, total_alerts)
    escalated_comparison = (previous_escalated, escalated_alerts)

    # Top alert types
    top_alerts = df["Alert Type"].value_counts().head(4)
    top_alerts_dict = top_alerts.to_dict()
    top_alerts_values = tuple(top_alerts.values)

    # Event sources (assuming a separate sheet or section; placeholder for now)
    event_sources = [
        ["Active Directory", "APA-SADC05-AD", "Running"],
        ["Cloud Service", "APA-M365-Logs", "Running"],
        # ... (17 rows as per March template, simplified here)
    ]

    # Accounts created/locked
    accounts_created = len(df[df["Alert Type"] == "New Local User Account Created"])
    accounts_locked = len(df[df["Alert Type"] == "User Behavior - A User Account Was Disabled"])
    accounts = {"created": accounts_created, "locked": accounts_locked}

    # Password resets
    password_resets = len(df[df["Alert Type"] == "User Behavior - An Attempt Was Made To Reset An Account's Password"])

    # Admin users and non-expiring accounts (placeholders)
    admin_users = "TBD"  # Needs actual data from Excel
    non_expiring = 2  # Placeholder from March

    return {
        "month_year": month_year,
        "month": month,
        "total_alerts": total_alerts,
        "escalated_alerts": escalated_alerts,
        "disposition_breakdown": (not_reported, reported),
        "priority_breakdown": priority_breakdown,
        "alerts_comparison": alerts_comparison,
        "escalated_comparison": escalated_comparison,
        "top_alerts": top_alerts_values,
        "top_alerts_categories": list(top_alerts_dict.keys()),
        "event_sources": event_sources,
        "accounts": accounts,
        "password_resets": password_resets,
        "admin_users": admin_users,
        "non_expiring": non_expiring
    }

# Add a title to a slide as a text box
def add_slide_title(slide, title_text, font_color=RGBColor(255, 255, 255)):
    # Add the title as a text box at the top of the slide
    tx_box = slide.shapes.add_textbox(Inches(3.5), Inches(1.2), Inches(8), Inches(1))
    tf = tx_box.text_frame
    p = tf.add_paragraph()
    p.text = title_text
    p.font.size = Pt(32)
    p.font.color.rgb = font_color  # Use the provided font color (default: dark blue)
    p.alignment = PP_ALIGN.CENTER

# Add a text box to a slide
def add_text_box(slide, text, left, top, width, height, font_color=RGBColor(0, 0, 0)):
    tx_box = slide.shapes.add_textbox(left, top, width, height)
    tf = tx_box.text_frame
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(20)
    p.font.color.rgb = font_color  # Use the provided font color (default: black)
    p.alignment = PP_ALIGN.CENTER

# Add a bar chart to a slide
def add_bar_chart(slide, chart_data, left, top, width, height, title):
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)

# Add a table to a slide
def add_table(slide, rows, cols, left, top, width, height, data):
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    for i in range(rows):
        for j in range(cols):
            table.cell(i, j).text = str(data[i][j])
            table.cell(i, j).text_frame.paragraphs[0].font.size = Pt(12)