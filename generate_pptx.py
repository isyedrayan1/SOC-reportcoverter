from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from config import SLIDE_CONFIG
from utils import add_slide_title, add_text_box, add_bar_chart, add_table

def generate_pptx(data, output_file):
    # Start with a blank presentation
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for slide_config in SLIDE_CONFIG:
        slide_number = slide_config["slide_number"]
        slide_type = slide_config["type"]
        title = slide_config["title"]

        # Format title with dynamic data if needed
        if "{month_year}" in title:
            title = title.format(month_year=data["month_year"])
        if "{month}" in title:
            title = title.format(month=data["month"])

        # Use a blank slide layout for all slides
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Set background image based on slide number
        if slide_number == 1:
            # First slide (img2)
            background_image = "based/img2.png"
        elif slide_number == 21:
            # Last slide (img3 - Thank You slide)
            background_image = "based/img3.png"
        else:
            # Middle slides (img1 - header/footer)
            background_image = "based/img1.png"

        # Add the background image
        slide.shapes.add_picture(
            background_image,
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height
        )

        # Add content based on slide number and type
        if slide_number == 1:
            # First slide: Add title and date on the left (both in white)
            #add_slide_title(slide, title, font_color=RGBColor(255, 255, 255))  # White title
            # Add the date on the left side
            add_text_box(slide, data["month_year"], Inches(1.4), Inches(4.3), Inches(3), Inches(1), font_color=RGBColor(255, 255, 255))  # White date

        elif slide_number == 21:
            # Last slide: No additional content needed (just the background image)
            pass

        else:
            # Middle slides (2-20): Add title and content on top of the background
            add_slide_title(slide, title)  # Default color (dark blue)

            # Add content based on slide type
            if slide_type == "title":
                add_text_box(slide, data["month_year"], Inches(3.665), Inches(2.5), Inches(6), Inches(2))  # Adjusted top to avoid header

            elif slide_type == "text":
                value = data[slide_config["data_key"]]
                if isinstance(value, dict):  # For accounts created/locked
                    text = slide_config["format"].format(**value)
                else:
                    text = slide_config["format"].format(value=value)
                add_text_box(slide, text, Inches(3.665), Inches(2.5), Inches(6), Inches(2))  # Adjusted top to avoid header

            elif slide_type == "bar_chart":
                chart_data = CategoryChartData()
                categories = slide_config["categories"]
                if categories == "dynamic":
                    categories = data[f"{slide_config['data_key']}_categories"]
                chart_data.categories = categories
                chart_data.add_series("Values", data[slide_config["data_key"]])
                add_bar_chart(slide, chart_data, Inches(2.665), Inches(2), Inches(8), Inches(4), slide_config["chart_title"])  # Adjusted top to avoid header

            elif slide_type == "table":
                table_data = [slide_config["columns"]] + data[slide_config["data_key"]]
                add_table(slide, len(table_data), len(slide_config["columns"]), Inches(3.665), Inches(2.5), Inches(8), Inches(4), table_data)  # Adjusted top to avoid header

            elif slide_type == "static_text":
                add_text_box(slide, slide_config["content"], Inches(3.665), Inches(2.5), Inches(6), Inches(2))  # Adjusted top to avoid header

    # Save the updated presentation
    prs.save(output_file)
    return output_file
